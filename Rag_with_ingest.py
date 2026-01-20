import ollama
import os
import uuid
import re
import pickle
import numpy as np
import chromadb
import nltk
import faiss
import pandas as pd
import cohere 
from posthog import page
import weaviate
import pytesseract
import torch


from email.mime import image, text
from csv import reader
from pptx import Presentation
from PIL import Image
from openpyxl import load_workbook
from docx import Document
from weaviate.auth import AuthApiKey
from weaviate.classes.config import Property
from weaviate.classes.config import DataType as WeaviateDataType
from weaviate.classes.data import DataObject
from sklearn.metrics.pairwise import cosine_similarity
from pypdf import PdfReader 
from pdf2image import convert_from_path
from nltk.tokenize import sent_tokenize
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings

from qdrant_client import QdrantClient
from qdrant_client.models import PointStruct, VectorParams, Distance
# from pymilvus import connections, FieldSchema, CollectionSchema, DataType, Collection,utility
from pinecone import Pinecone, ServerlessSpec

from transformers import BlipProcessor, BlipForConditionalGeneration     #BlipProcessor   , BlipForConditionalGeneration


# ==============================
# INITIAL SETUP
# ==============================

nltk.download("punkt")
nltk.download("punkt_tab")
CHROMA_COLLECTION = "rag-chunks"
LLM_MODEL = "llama3:8b"
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CHROMA_PATH = os.path.join(BASE_DIR, "chroma_db")       # ----- chroma path
FAISS_PATH = os.path.join(BASE_DIR, "faiss_store.pkl")     # ----- faiss path
QDRANT_PATH = os.path.join(BASE_DIR, "qdrant_db")       # ----- qdrant path
QDRANT_COLLECTION = "rag_chunks"                          # ----- qdrant collection name
MILVUS_COLLECTION = "rag_chunks"                           # ----- milvus collection name
PINECONE_API_KEY = "pcsk_36kLSk_KyRnADaFcY1xF4dS7gyTKjCqBu5qHj8rYMmmr7RP6anZb56YmGZw1Z5BaM1Gznv"      # pinecone api key
PINECONE_ENV = "us-east-1"                                                                   # ---------------- pinecone environment
PINECONE_INDEX = "rag-chunks"                                            
WEAVIATE_URL = "https://pmh8cdktsmeloccbguom5a.c0.asia-southeast1.gcp.weaviate.cloud"                  # ---------------- weaviate url
WEAVIATE_API_KEY = "UGd4eU0yeC9sbWhNZm53Nl9STDIvSTY0K2puOTRUY05KZ1hNVEhYajhIbk9PTUVMKzluc3RERVI0WC9NPV92MjAw"   # ---------------- weaviate api key
WEAVIATE_CLASS = "RagChunks"
COHERE_API_KEY="BALyrvraWjAJ0273i1sNKyGNs4S2ase0Y8I2ZezG" #------------------- cohere api key


DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")          #------ "Salesforce/blip-image-captioning-base" for cpu
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base").to(DEVICE)  # -----"Salesforce/blip-image-captioning-base" for cpu

# ==============================
# AGENT RULES (CONFIG)
# ==============================
    
AGENT_RULES = {
    "technical_assistant": {
        "keep_keywords": [
            "architecture", "design", "deployment", "installation",
            "configuration", "setup", "implementation"
        ],
        "drop_keywords": [
            "pricing", "contact", "about", "company", "legal", "copyright"
        ],
        "merge_keywords": [
            "setup", "installation", "deployment"
        ],
        "max_section_length": 1500
    }
}

# ==============================
# EMBEDDING REGISTRY (PRODUCTION)
# ==============================
EMBEDDING_REGISTRY = {
    "bge-base": {
        "model_name": "BAAI/bge-base-en-v1.5",
        "provider": "huggingface",
        "type": "default"
    },
    "bge-m3": {
        "model_name": "BAAI/bge-m3",
        "provider": "huggingface",
        "type": "multilingual"
    },
    "e5-large": {
        "model_name": "intfloat/e5-large-v2",
        "provider": "huggingface",
        "type": "high_accuracy",
        "requires_prefix": True
    },
    "e5-multilingual": {
        "model_name": "intfloat/multilingual-e5-large",
        "provider": "huggingface",
        "type": "multilingual",
        "requires_prefix": True
    },
    "gte-large": {
        "model_name": "thenlper/gte-large",
        "provider": "huggingface",
        "type": "instruction"
    },
    "jina-v3": {
        "model_name": "jinaai/jina-embeddings-v3",
        "provider": "huggingface",
        "type": "infra_scale"
    },
    "nomic-v1.5": {
        "model_name": "nomic-ai/nomic-embed-text-v1.5",
        "provider": "huggingface",
        "type": "open_source"
    },

    # ✅ COHERE (API-based, Windows-safe)
    "cohere-english": {
        "model_name": "embed-english-v3.0",
        "provider": "cohere",
        "type": "api"
    },
    "cohere-multilingual": {
        "model_name": "embed-multilingual-v3.0",
        "provider": "cohere",
        "type": "api"
    }
}


# ==============================
# DOCUMENT LOADER CLASS (UNIFIED)
# ==============================

class DocumentLoader:
    """
    Unified document loader for RAG ingestion.
    Returns a list of {text, metadata} dicts.
    """

    def load(self, file_path: str):
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self._load_pdf(file_path)

        elif ext in [".txt", ".text", ".log"]:
            return self._load_text(file_path)

        elif ext in [".md", ".markdown", ".mdown", ".mkd"]:
            return self._load_markdown(file_path)

        elif ext in [".docx"]:
            return self._load_docx(file_path)

        elif ext in [".csv", ".tsv", ".xls", ".xlsx", ".xlsm"]:
            return self._load_table(file_path)

        elif ext in [".pptx"]:
            return self._load_pptx(file_path)

        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"]:
            return self._load_image(file_path)

        else:
            raise ValueError(f"Unsupported file type: {ext}")

    # ---------- LOADERS ----------

    def _load_pdf(self, file_path):      #----------------------------PDF LOADER WITH OCR + IMAGE CAPTIONING
        docs = []

    # ==============================
    # 1️⃣ Normal text extraction
    # ==============================
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                docs.append({
                    "text": text.strip(),
                    "metadata": {
                    "source": os.path.basename(file_path),
                    "page": i + 1,
                    "type": "pdf_text",
                    "confidence": "high"
                }
            })

    # ==============================
    # 2️⃣ OCR + Image Captioning
    # ==============================
        try:
            images = convert_from_path(file_path)

            for i, img in enumerate(images):
            # ---- OCR ----
              ocr_text = pytesseract.image_to_string(img)
              if ocr_text and ocr_text.strip():
                docs.append({
                    "text": ocr_text.strip(),
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": i + 1,
                        "type": "pdf_ocr",
                        "confidence": "high"
                    }
                })

            # ---- Image captioning ----
            caption = self.describe_image(img)
            if caption and len(caption.split()) >= 4:
                docs.append({
                    "text": caption,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": i + 1,
                        "type": "pdf_image_caption",
                        "confidence": "medium"
                    }
                })

        except Exception as e:
            print("⚠ OCR / Image captioning skipped:", e)

        return docs

#=======================================================================================

    def _load_text(self, file_path):         # ----------------------------TEXT LOADER
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        return [{
            "text": text,
            "metadata": {
                "source": os.path.basename(file_path),
                "page": 1,
                "type": "text"
            }
        }]
#========================================================================================

    def _load_markdown(self, file_path):             # ----------------------------MARKDOWN LOADER
        # Markdown is treated as plain text for RAG
        return self._load_text(file_path)
    

#========================================================================================
    def _load_docx(self, file_path):                    # ----------------------------DOCX LOADER
        doc = Document(file_path)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        return [{
            "text": text,
            "metadata": {
                "source": os.path.basename(file_path),
                "page": 1,
                "type": "docx"
            }
        }]

    def _load_table(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".csv":
            df = pd.read_csv(file_path)
        elif ext == ".tsv":
            df = pd.read_csv(file_path, sep="\t")
        else:
            df = pd.read_excel(file_path)

        df = df.dropna(how="all").fillna("")
        docs = []

        for idx, row in df.iterrows():
            row_text = " | ".join(
                f"{col}: {str(row[col]).strip()}"
                for col in df.columns
                if str(row[col]).strip()
            )

            if row_text:
                docs.append({
                    "text": row_text,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "row": int(idx),
                        "type": "table"
                    }
                })

        return docs
#========================================================================================

    def _load_pptx(self, file_path):      # ----------------------------PPTX LOADER
        prs = Presentation(file_path)
        docs = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            if slide_text:
                docs.append({
                    "text": "\n".join(slide_text),
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": i + 1,
                        "type": "pptx"
                    }
                })

        return docs
    
#========================================================================================
    # ==============================
    # IMAGE LOADER (OCR + CAPTION)
    # ==============================
    def _load_image(self, image_path):                     #----------------------------IMAGE LOADER WITH OCR + IMAGE CAPTIONING
        docs = []

        try:
            image = Image.open(image_path).convert("RGB")
        except Exception as e:
            print(f"⚠ Failed to open image {image_path}: {e}")
            return docs

        # 1️⃣ OCR
        ocr_text = pytesseract.image_to_string(image)
        if ocr_text.strip():
            docs.append({
                "text": ocr_text.strip(),
                "metadata": {
                    "source": os.path.basename(image_path),
                    "page": 1,
                    "type": "image_ocr",
                    "confidence": "high"
                }
            })

        # 2️⃣ Image captioning
        caption = self.describe_image(image)
        if caption and len(caption.split()) >= 4:
            docs.append({
                "text": caption,
                "metadata": {
                    "source": os.path.basename(image_path),
                    "page": 1,
                    "type": "image_caption",
                    "confidence": "medium"
                }
            })

        return docs

    # ==============================
    # IMAGE CAPTIONING
    # ==============================
    def describe_image(self, image: Image.Image) -> str:
        try:
            inputs = blip_processor(
                image,
                return_tensors="pt"
            ).to(DEVICE)
            with torch.no_grad():
                 output = blip_model.generate(
                     **inputs,
                    max_length=60,
                    num_beams=5
                )

            caption = blip_processor.decode(
                output[0],
                skip_special_tokens=True
            )
            return caption.strip()
        except Exception as e:
             print("⚠ Image captioning failed:", e)
        return ""

# ==============================
# BASIC CHUNKING STRATEGIES
# ==============================

def fixed_chunking(text, chunk_size=500, overlap=50):   #----------------- fixed chunking function
    chunks, start = [], 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks


def sentence_chunking(text, max_sentences=5):    #------------------ sentence chunking function
    sentences = sent_tokenize(text, language="english")
    return [
        " ".join(sentences[i:i + max_sentences])
        for i in range(0, len(sentences), max_sentences)
    ]


def paragraph_chunking(text):              #------------------ paragraph chunking function
    return [p.strip() for p in text.split("\n\n") if p.strip()]


def recursive_chunking(text, chunk_size=500, overlap=50):    #------------------ recursive chunking function
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=overlap
    )
    return splitter.split_text(text)

# ==============================
# SEMANTIC CHUNKING
# ==============================

def semantic_chunking(
    text,
    embedder,
    similarity_threshold=0.75,
    max_chunk_chars=800
):                                                             #------------------ semantic chunking function
    sentences = sent_tokenize(text, language="english")
    if len(sentences) <= 1:
        return sentences

    sentence_embeddings = embedder.embed_documents(sentences)

    chunks = []
    current_chunk = [sentences[0]]
    current_len = len(sentences[0])

    for i in range(1, len(sentences)):
        sim = cosine_similarity(
            [sentence_embeddings[i - 1]],
            [sentence_embeddings[i]]
        )[0][0]

        if sim >= similarity_threshold and current_len < max_chunk_chars:
            current_chunk.append(sentences[i])
            current_len += len(sentences[i])
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [sentences[i]]
            current_len = len(sentences[i])

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

# ==============================
# CONTENT-AWARE CHUNKING
# ============================== 

def is_heading(line: str) -> bool:
    line = line.strip()
    if not line:
        return False
    if line.isupper() and len(line) < 100:
        return True
    if re.match(r"^\d+(\.\d+)*\s+", line):
        return True
    if line.endswith(":") and len(line) < 80:
        return True
    return False


def detect_sections(text: str):
    lines = text.split("\n")
    sections = []

    current = {
        "title": "INTRODUCTION",
        "content": [],
        "index": 0
    }

    section_index = 1

    for line in lines:
        if is_heading(line):
            if current["content"]:
                sections.append(current)

            current = {
                "title": line.strip(),
                "content": [],
                "index": section_index
            }
            section_index += 1
        else:
            current["content"].append(line)

    if current["content"]:
        sections.append(current)

    return sections


def content_aware_chunking(text, embedder, inner_strategy="recursive"):
    sections = detect_sections(text)
    chunks = []

    for section in sections:
        section_text = "\n".join(section["content"]).strip()
        if not section_text:
            continue

        if inner_strategy == "paragraph":
            inner_chunks = paragraph_chunking(section_text)
        elif inner_strategy == "semantic":
            inner_chunks = semantic_chunking(section_text, embedder)
        else:
            inner_chunks = recursive_chunking(section_text)

        for idx, chunk in enumerate(inner_chunks):
            chunks.append({
                "text": chunk,
                "section_title": section["title"],
                "section_index": section["index"],
                "chunk_index_in_section": idx
            })

    return chunks

# ==============================
# AGENTIC CHUNKING
# ==============================

def evaluate_section(section, agent_goal):
    rules = AGENT_RULES[agent_goal]
    title = section["title"].lower()
    content = " ".join(section["content"]).lower()

    for kw in rules["drop_keywords"]:
        if kw in title:
            return "DROP", f"Irrelevant section: {section['title']}"

    for kw in rules["merge_keywords"]:
        if kw in title:
            return "MERGE", f"Merged operational section: {section['title']}"

    if len(content) > rules["max_section_length"]:
        return "SPLIT", "Section too large"

    for kw in rules["keep_keywords"]:
        if kw in title or kw in content:
            return "KEEP", "Relevant to agent goal"

    return "KEEP", "Default keep"


def agentic_chunking(
    text,
    embedder,
    agent_goal="technical_assistant",
    inner_strategy="recursive"
):
    sections = detect_sections(text)
    chunks = []
    merge_buffer = []

    for section in sections:
        decision, reason = evaluate_section(section, agent_goal)
        section_text = "\n".join(section["content"]).strip()

        if not section_text:
            continue

        if decision == "DROP":
            continue

        if decision == "MERGE":
            merge_buffer.append(section)
            continue

        if merge_buffer:
            merged_text = "\n\n".join(
                ["\n".join(s["content"]) for s in merge_buffer]
            )
            chunks.append({
                "text": merged_text,
                "decision": "MERGE",
                "decision_reason": "Merged related sections",
                "agent_goal": agent_goal,
                "source_sections": [s["title"] for s in merge_buffer]
            })
            merge_buffer = []

        if decision == "SPLIT":
            sub_chunks = (
                semantic_chunking(section_text, embedder)
                if inner_strategy == "semantic"
                else recursive_chunking(section_text)
            )
            for idx, sub in enumerate(sub_chunks):
                chunks.append({
                    "text": sub,
                    "decision": "SPLIT",
                    "decision_reason": reason,
                    "agent_goal": agent_goal,
                    "source_sections": [section["title"]],
                    "chunk_index_in_section": idx
                })

        if decision == "KEEP":
            chunks.append({
                "text": section_text,
                "decision": "KEEP",
                "decision_reason": reason,
                "agent_goal": agent_goal,
                "source_sections": [section["title"]]
            })

    return chunks

# ==============================
# EMBEDDING MODEL
# ==============================
                        
class CohereEmbeddingWrapper:                                        #------------------ cohere embedding wrapper class
    def __init__(self, model_name: str):
        api_key = os.getenv("COHERE_API_KEY") or COHERE_API_KEY
        if not api_key:
            raise ValueError("COHERE_API_KEY not found")
        self.client = cohere.Client(api_key)
        self.model_name = model_name

    def embed_documents(self, texts):
        response = self.client.embed(
            model=self.model_name,
            texts=texts,
            input_type="search_document"
        )
        return response.embeddings

    def embed_query(self, text):
        response = self.client.embed(
            model=self.model_name,
            texts=[text],
            input_type="search_query"
        )
        return response.embeddings[0]
    def close(self):
        if hasattr(self.client, "close"):
            self.client.close()
    

def load_embedding_model(embedding_key: str):
    if embedding_key not in EMBEDDING_REGISTRY:
        raise ValueError(f"Unsupported embedding model: {embedding_key}")

    config = EMBEDDING_REGISTRY[embedding_key]
    
    if config["provider"] == "huggingface":
        return HuggingFaceEmbeddings(
            model_name=config["model_name"]
        )
    elif config["provider"] == "cohere":
        return CohereEmbeddingWrapper(
            model_name=config["model_name"]
        )
    else :
      raise ValueError("Unsupported embedding provider")

# ==============================
# helper function for E5
# ==============================
def prepare_texts_for_embedding(texts, embedding_key, is_query=False):
    config = EMBEDDING_REGISTRY[embedding_key]

    if config.get("requires_prefix"):
        prefix = "query: " if is_query else "passage: "
        return [prefix + t for t in texts]

    return texts


# ==============================
# VECTOR STORE ADAPTERS
# ==============================

class VectorStore:
    def add(self, texts, embeddings, metadatas, ids):
        raise NotImplementedError
    
    def search(self, query_embedding, top_k=5):
        raise NotImplementedError

    def count(self):
        raise NotImplementedError


class ChromaStore(VectorStore):
    def __init__(self,suffix):
        os.makedirs(CHROMA_PATH, exist_ok=True)
        self.client = chromadb.PersistentClient(path=CHROMA_PATH)
        self.collection = self.client.get_or_create_collection(name=f"rag_chunks__{suffix}")

    def add(self, texts, embeddings, metadatas, ids):
        self.collection.add(
            documents=texts,
            embeddings=embeddings,
            metadatas=metadatas,
            ids=ids
        )
    def search(self, query_embedding, top_k=5):
        res = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )
        return res

    def count(self):
        return self.collection.count()


class FaissStore(VectorStore):
    def __init__(self, dim, suffix):
        self.index = faiss.IndexFlatIP(dim)
        self.path = os.path.join(BASE_DIR, f"faiss_store__{suffix}.pkl")
        self.data = []

        if os.path.exists(self.path):
            with open(self.path, "rb") as f:
                saved = pickle.load(f)
                self.index = saved["index"]
                self.data = saved["data"]

    def add(self, texts, embeddings, metadatas, ids): 
        vectors = np.array(embeddings).astype("float32")
        self.index.add(vectors)

        for i in range(len(texts)):
            self.data.append({
                "id": ids[i],
                "text": texts[i],
                "metadata": metadatas[i]
            })

        with open(self.path, "wb") as f:
            pickle.dump({"index": self.index, "data": self.data}, f)

    def search(self, query_embedding, top_k=5):
        query_vector = np.array([query_embedding]).astype("float32")
        distances, indices = self.index.search(query_vector, top_k)

        results = []
        for dist, idx in zip(distances[0], indices[0]):
            if idx < len(self.data):
                item = self.data[idx]
                item["score"] = float(dist)
                results.append(item)
        return results

    def count(self):
        return self.index.ntotal


class QdrantStore(VectorStore):
    def __init__(self, dim, suffix):
        self.client = QdrantClient(url="https://ec558055-1e28-4fa3-9999-92fd632a068f.europe-west3-0.gcp.cloud.qdrant.io:6333", 
                                   api_key="eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJhY2Nlc3MiOiJtIn0.XXeIwZZ4e1sa5UbH6Bix9t0_HWL1N3VD0XrKxM0HnfE",
                                   )
        self.collection_name = f"rag_chunks__{suffix}"

        existing = [c.name for c in self.client.get_collections().collections]
        if self.collection_name not in existing:
            self.client.create_collection(
                collection_name=self.collection_name,
                vectors_config=VectorParams(
                    size=dim,
                    distance=Distance.COSINE
                )
            )

    def add(self, texts, embeddings, metadatas, ids):
        points = []
        for i in range(len(texts)):
            points.append(
                PointStruct(
                    id=ids[i],
                    vector=embeddings[i],
                    payload={
                        "text": texts[i],
                        **metadatas[i]
                    }
                )
            )

        self.client.upsert(
            collection_name=self.collection_name,
            points=points
        )

    def search(self, query_embedding, top_k=5):
         return self.client.query_points(
                collection_name=self.collection_name,
                query=query_embedding,
                limit=top_k,
                with_payload=True,  
                with_vectors=False
            )

    def count(self):
        return self.client.count(self.collection_name).count


    def close(self):
        self.client.close()
# # ==============================
# # ✅ MILVUS STORE (ADDED)
# # ==============================

# class MilvusStore(VectorStore):
#     def __init__(self, dim,suffix):
#         connections.connect(
#             alias="default",
#             host="localhost",
#             port="19530"
#         )

#         if not utility.has_collection(MILVUS_COLLECTION):
#             fields = [
#                 FieldSchema(
#                     name="id",
#                     dtype=DataType.VARCHAR,
#                     is_primary=True,
#                     max_length=64
#                 ),
#                 FieldSchema(
#                     name="embedding",
#                     dtype=DataType.FLOAT_VECTOR,
#                     dim=dim
#                 ),
#                 FieldSchema(
#                     name="text",
#                     dtype=DataType.VARCHAR,
#                     max_length=65535
#                 ),
#                 FieldSchema(
#                     name="metadata",
#                     dtype=DataType.JSON
#                 )
#             ]

#             schema = CollectionSchema(
#                 fields=fields,
#                 description="RAG chunks"
#             )

#             self.collection = Collection(
#                 name=MILVUS_COLLECTION,
#                 schema=schema
#             )

#             self.collection.create_index(
#                 field_name="embedding",
#                 index_params={
#                     "index_type": "HNSW",
#                     "metric_type": "COSINE",
#                     "params": {"M": 8, "efConstruction": 64}
#                 }
#             )
#         else:
#             self.collection = Collection(MILVUS_COLLECTION)

#         self.collection.load()

#     def add(self, texts, embeddings, metadatas, ids):
#         self.collection.insert([
#             ids,
#             embeddings,
#             texts,
#             metadatas
#         ])

#     def count(self):
#         return self.collection.num_entities

#     def close(self):
#         self.collection.release()
#         connections.disconnect("default")


# ==============================
# ✅ Pine-cone
# ==============================
class PineconeStore(VectorStore):
    def __init__(self, dim, suffix):
        self.pc = Pinecone(api_key=PINECONE_API_KEY)
        index_name = f"rag-chunks-{suffix}"

        existing_indexes = [idx["name"] for idx in self.pc.list_indexes()]
        if index_name not in existing_indexes:
            self.pc.create_index(
                name=index_name,
                dimension=dim,
                metric="cosine",
                spec=ServerlessSpec(
                    cloud="aws",
                    region=PINECONE_ENV
                )
            )

        self.index = self.pc.Index(index_name)

    def add(self, texts, embeddings, metadatas, ids):
        vectors = []
        for i in range(len(texts)):
            vectors.append((
                ids[i],
                embeddings[i],
                {
                    "text": texts[i],
                    **metadatas[i]
                }
            ))

        self.index.upsert(vectors=vectors)
    
    def search(self, query_embedding, top_k=5):
        res = self.index.query(
            vector=query_embedding,
            top_k=top_k,
            include_metadata=True
        )
        return res["matches"]

    def count(self):
        stats = self.index.describe_index_stats()
        return stats["total_vector_count"]
    


# ==============================
# ✅ Weaviate STORE
# ==============================
class WeaviateStore(VectorStore):
    def __init__(self, dim, suffix):
        # 1️⃣ Connect to Weaviate Cloud (v4)
        self.client = weaviate.connect_to_weaviate_cloud(
            cluster_url=WEAVIATE_URL.replace("https://", ""),
            auth_credentials=AuthApiKey(WEAVIATE_API_KEY),
        )
        class_name = f"RagChunks_{suffix.replace('-', '_')}"
        print("🔍 [Weaviate] Client created")

        ready = self.client.is_ready()
        print("🔍 [Weaviate] is_ready():", ready)

        if not ready:
            raise RuntimeError("Weaviate is not ready")
        existing = self.client.collections.list_all()      
        if class_name not in existing:
            self.client.collections.create(
                      name=class_name,
                         properties=[
                            Property(name="text", data_type=WeaviateDataType.TEXT),
                            Property(name="source", data_type=WeaviateDataType.TEXT),
                            Property(name="page", data_type=WeaviateDataType.INT),
                        ],
                        )
        else:
            print(f"✅ [Weaviate] Collection already exists: {class_name}")
        # 5️⃣ Get collection handle
        self.collection = self.client.collections.get(class_name)

    def add(self, texts, embeddings, metadatas, ids):
        
        objects = []
        for i in range(len(texts)):
             objects.append(
            DataObject(
                properties={
                    "text": texts[i],
                    **metadatas[i]
                },
                vector=embeddings[i],
                uuid=ids[i]
            )
        )
        self.collection.data.insert_many(objects)

    def search(self, query_embedding, top_k=5):
        res = self.collection.query.near_vector(
            near_vector=query_embedding,
            limit=top_k
        )
        return res.objects
    
    def count(self):
        count = self.collection.aggregate.over_all(
            total_count=True
        ).total_count

        print("🔢 [Weaviate] Total vectors:", count)
        return count
    
    def close(self):
        print("🔒 [Weaviate] Closing client")
        self.client.close()



def init_vector_store(db_type, embedder, embedding_key):
    sample_vec = embedder.embed_query("dimension test")
    dim = len(sample_vec)
    suffix = embedding_key
    if db_type == "chroma":
        return ChromaStore(suffix)
    elif db_type == "faiss":
        return FaissStore(dim, suffix)
    elif db_type == "qdrant":
        return QdrantStore(dim, suffix)
    # elif db_type == "milvus":          
    #     return MilvusStore(dim,suffix)
    elif db_type == "pinecone":      
        return PineconeStore(dim, suffix)
    elif db_type == "weaviate":
        return WeaviateStore(dim, suffix)
    else:
        raise ValueError("Unsupported DB")
    
# ==============================
# Rag-Model PIPELINE
# ==============================

def generate_answer(query, context):
    system_prompt = f"""
You are a document-grounded assistant.

Rules:
- Use ONLY the provided context
- Do NOT hallucinate
- If answer not found, say: "Not found in the document"

Context:
{context}
"""
    response = ollama.chat(
        model=LLM_MODEL,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": query}
        ]
    )
    return response["message"]["content"]
 
def embed_query_text(query, embedder, embedding_key):
    query = prepare_texts_for_embedding(
        [query],
        embedding_key,
        is_query=True
    )[0]
    return embedder.embed_query(query)


def retrieve_top_k(
    query,
    embedder,
    vector_store,
    embedding_key,
    top_k=5
):
    query_embedding = embed_query_text(
        query,
        embedder,
        embedding_key
    )
    results = vector_store.search(
        query_embedding,
        top_k=top_k
    )
    return results


def build_context(results, vector_db):
    contexts = []


    if vector_db == "qdrant":
        for r in results:
            # r is a tuple → unpack it safely
            if len(r) == 3:
                _, score, payload = r
            elif len(r) == 2:
                _, score = r
                continue
            else:
                continue

            if score > 0.2 and "text" in payload:
                contexts.append(payload["text"])

    elif vector_db == "faiss":
        for r in results:
            contexts.append(r["text"])

    elif vector_db == "chroma":
        # Chroma returns a single dict, NOT per-item results
        if "documents" in results and len(results["documents"]) > 0:
            contexts.extend(results["documents"][0])

    elif vector_db == "pinecone":
        for r in results:
            contexts.append(r["metadata"]["text"])

    elif vector_db == "weaviate":
        for r in results:
            contexts.append(r.properties["text"])

    return "\n\n".join(contexts)

def rag_query(
    query,
    embedder,
    vector_store,
    embedding_key,
    vector_db,
    top_k=5
):
    results = retrieve_top_k(
        query,
        embedder,
        vector_store,
        embedding_key,
        top_k
    )

    MAX_CONTEXT_CHARS = 6000
    context = build_context(results, vector_db)
    context = context[:MAX_CONTEXT_CHARS]
    answer = generate_answer(query, context)

    return {
        "query": query,
        "answer": answer,
        "context": context,
        "sources": results
    }

# ==============================
# INGESTION PIPELINE
# ==============================

def ingest_document(file_path, strategy, collection, embedder, embedding_key):
    loader = DocumentLoader()
    documents = loader.load(file_path)
    total_chunks = 0

    for doc in documents:
        text = doc["text"]
        base_metadata = doc["metadata"]

        if strategy == "agentic":
            chunks = agentic_chunking(text, embedder)
            texts = [c["text"] for c in chunks]
        elif strategy == "content_aware":
            chunks = content_aware_chunking(text, embedder)
            texts = [c["text"] for c in chunks]
        elif strategy == "semantic":
            chunks = semantic_chunking(text, embedder)
            texts = chunks
        elif strategy == "sentence":
            chunks = sentence_chunking(text)
            texts = chunks
        elif strategy == "paragraph":
            chunks = paragraph_chunking(text)
            texts = chunks
        elif strategy == "fixed":
            chunks = fixed_chunking(text)
            texts = chunks
        else:
            chunks = recursive_chunking(text)
            texts = chunks

        prepared_texts = prepare_texts_for_embedding(
             texts,
             embedding_key,
            is_query=False
            )

        embeddings = embedder.embed_documents(prepared_texts)
        for i, emb in enumerate(embeddings):
            if strategy == "agentic":
                meta = {
                    **base_metadata,
                    "chunk_type": "agentic",
                    "agent_goal": chunks[i]["agent_goal"],
                    "decision": chunks[i]["decision"],
                    "decision_reason": chunks[i]["decision_reason"],
                    "source_sections": " | ".join(chunks[i]["source_sections"])
                }
                chunk_text = chunks[i]["text"]
            elif strategy == "content_aware":
                meta = {
                    **base_metadata,
                    "chunk_type": "content_aware",
                    "section_title": chunks[i]["section_title"],
                    "section_index": chunks[i]["section_index"],
                    "chunk_index_in_section": chunks[i]["chunk_index_in_section"]
                }
                chunk_text = chunks[i]["text"]
            else:
                meta = {**base_metadata, "chunk_type": strategy}
                chunk_text = texts[i]

            collection.add(
                 texts=[chunk_text],   
                embeddings=[emb],
                metadatas=[meta],
                ids=[str(uuid.uuid4())]
            )       
            total_chunks += 1

    return total_chunks

# ==============================
# ENTRY POINT
# ==============================

if __name__ == "__main__":

    pdf_path = "ai demo.jpg"
    strategy = "agentic"          # fixed | sentence | paragraph | recursive | semantic | agentic | content_aware
    vector_db = "qdrant"            # chroma | faiss | qdrant | milvus | pinecone | weaviate
    embedding_key = "bge-base"      # bge-base | bge-m3 | e5-large | e5-multilingual | gte-large | nomic-v1.5 | cohere-english | cohere-multilingual
    print("🧠 Loading embedding model...")
    embedder = load_embedding_model(embedding_key)
    print(f"🧠 Embedding model loaded: {embedding_key}")
    print(f"🗄️ Initializing Vector DB: {vector_db}")
    vector_store = init_vector_store(vector_db, embedder,embedding_key)

    print(f"📄 Ingesting document using '{strategy}' chunking...")
    total = ingest_document(pdf_path, strategy, vector_store, embedder,embedding_key)

    print("\n✅ INGESTION COMPLETE")
    print(f"📄 PDF: {pdf_path}")
    print(f"✂️ Chunking Strategy: {strategy}")
    print(f"🗄️ Vector DB: {vector_db}")
    print(f"📦 Total Chunks Stored: {total}")
    print("🔢 Total vectors in DB:", vector_store.count())
    print("\n🤖 RAG Chat Ready (Ctrl+C to exit)\n")
    while True:
        try:
            q = input("❓ Question: ").strip()
            if not q:
               continue

            result = rag_query(
                q,
                embedder,
                vector_store,
                embedding_key,
            vector_db
        )

            print("\n🧠 Answer:\n", result["answer"], "\n")
        except KeyboardInterrupt:
             break     
    print("🔒 Closing vector store...")  
    if hasattr(vector_store, "close"):
         vector_store.close() 